#!python3
"""
Office file operations via Microsoft Graph API + python-docx/python-pptx.

Excel  — Graph Workbook API (native read/write, no download needed)
Word   — download .docx, parse with python-docx, upload back
PowerPoint — download .pptx, parse with python-pptx, upload back

Requires:
    ~/.claude/skills/onenote/scripts/onenote_setup.py  (auth + Graph client)
    pip: python-docx, python-pptx
"""

import io
import json
import re
import sys
import argparse
import time
import warnings
import zipfile
import urllib.request
import urllib.parse
import urllib.error
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from datetime import datetime
warnings.filterwarnings('ignore', category=Warning, module='urllib3')

sys.path.insert(0, str(Path(__file__).parent.parent / 'onenote' / 'scripts'))
from onenote_setup import get_access_token

EXCEL_CACHE_DIR = Path.home() / '.claude' / 'skills' / 'office' / 'cache' / 'excel'
EXCEL_CACHE_DIR.mkdir(parents=True, exist_ok=True)

def _cache_path(file_id: str) -> Path:
    safe = file_id.replace('!', '_').replace('/', '_')
    return EXCEL_CACHE_DIR / f'{safe}.json'

def load_excel_cache(file_id: str):
    """Load cached tab descriptions for a file. Returns None if not cached.
    Auto-migrates legacy .md cache to .json on first access."""
    p = _cache_path(file_id)
    if not p.exists():
        # Migrate legacy .md if present
        md = p.with_suffix('.md')
        if md.exists():
            try:
                data = json.loads(md.read_text().split('```json\n')[1].split('\n```')[0])
                save_excel_cache(data['file_id'], data['filename'], data['tabs'])
                md.unlink()
            except Exception:
                pass
        if not p.exists():
            return None
    return json.loads(p.read_text())

def save_excel_cache(file_id: str, filename: str, tabs: dict,
                     file_modified: str = '', sheet_ids: dict[str, str] = None,
                     tab_groups: dict = None) -> None:
    """Save tab descriptions as compact JSON.
    tabs       = {sheet_name: str | list[{section, range, desc}]}
                 Use a string for simple single-section tabs.
                 Use a list of section dicts for multi-section/dashboard tabs.
    tab_groups = {group_name: {pattern, range, count, desc}} — covers groups of
                 structurally identical tabs (e.g. monthly expense sheets) so
                 they don't need individual entries in tabs.
    sheet_ids  = {sheet_name: graph_worksheet_id}  — used for rename detection
    file_modified = lastModifiedDateTime from Graph — used for staleness check"""
    p = _cache_path(file_id)
    doc = {'file_id': file_id, 'filename': filename,
           'updated': datetime.now().strftime('%Y-%m-%d'),
           'file_modified': file_modified,
           'sheet_ids': sheet_ids or {},
           'tabs': tabs}
    if tab_groups:
        doc['tab_groups'] = tab_groups
    p.write_text(json.dumps(doc, indent=2))

def is_excel_cache_stale(file_id: str, drive_id: str = None):
    """Check if the Excel file has been modified since the cache was built.
    Returns (is_stale: bool, current_modified: str).
    If no cache or no stored file_modified, returns (True, current_modified)."""
    cache = load_excel_cache(file_id)
    base = _item_path(file_id, drive_id)
    try:
        meta = _graph(f'{base}?$select=lastModifiedDateTime,name')
        current_mod = meta.get('lastModifiedDateTime', '')
    except Exception:
        return True, ''
    if not cache or not cache.get('file_modified'):
        return True, current_mod
    return cache['file_modified'] != current_mod, current_mod

# ---------------------------------------------------------------------------
# Drive helpers — find files on OneDrive
# ---------------------------------------------------------------------------

def _graph(path: str, method='GET', body=None, content_type='application/json') -> dict:
    token = get_access_token()
    url = f'https://graph.microsoft.com/v1.0{path}'
    data = json.dumps(body).encode() if body else None
    headers = {'Authorization': f'Bearer {token}'}
    if data:
        headers['Content-Type'] = content_type
    req = urllib.request.Request(url, data=data, headers=headers, method=method)
    with urllib.request.urlopen(req) as r:
        return json.loads(r.read()) if r.length != 0 else {}


def _graph_bytes(path: str) -> bytes:
    token = get_access_token()
    url = f'https://graph.microsoft.com/v1.0{path}'
    req = urllib.request.Request(url, headers={'Authorization': f'Bearer {token}'})
    with urllib.request.urlopen(req) as r:
        return r.read()


def _graph_upload(path: str, data: bytes, content_type: str) -> dict:
    token = get_access_token()
    url = f'https://graph.microsoft.com/v1.0{path}'
    req = urllib.request.Request(url, data=data,
                                  headers={'Authorization': f'Bearer {token}',
                                           'Content-Type': content_type},
                                  method='PUT')
    with urllib.request.urlopen(req) as r:
        return json.loads(r.read())


def search_files(query: str, file_type: str = None) -> list[dict]:
    """Search OneDrive for files matching query. Optionally filter by extension."""
    encoded = urllib.parse.quote(query)
    data = _graph(f'/me/drive/root/search(q=\'{encoded}\')?$select=id,name,webUrl,size,lastModifiedDateTime')
    results = data.get('value', [])
    if file_type:
        ext = file_type.lower().lstrip('.')
        results = [f for f in results if f['name'].lower().endswith(f'.{ext}')]
    return results


def list_folder(folder_path: str = '/') -> list[dict]:
    """List files in a OneDrive folder path."""
    if folder_path in ('/', ''):
        data = _graph('/me/drive/root/children?$select=id,name,webUrl,file,folder')
    else:
        encoded = urllib.parse.quote(folder_path)
        data = _graph(f'/me/drive/root:/{encoded}:/children?$select=id,name,webUrl,file,folder')
    return data.get('value', [])


def get_file_id(file_path: str) -> str:
    """Get Drive item ID for a file by path (e.g. 'Documents/report.xlsx')."""
    encoded = urllib.parse.quote(file_path)
    data = _graph(f'/me/drive/root:/{encoded}?$select=id,name')
    return data['id']


# ---------------------------------------------------------------------------
# Excel — session cache (avoids ~2.5s workbook-open overhead per request)
# ---------------------------------------------------------------------------

_excel_sessions: dict = {}   # (file_id, persist) -> {'id': str, 'ts': float}
SESSION_TTL = 240            # seconds — expire before Graph's ~5min timeout

def _get_session(file_id: str, drive_id: str = None, persist: bool = False) -> str:
    """Return a live workbook session ID, creating one if needed.
    persist=False for reads; persist=True for writes (commits changes to file)."""
    key = (file_id, persist)
    entry = _excel_sessions.get(key)
    if entry and (time.time() - entry['ts']) < SESSION_TTL:
        return entry['id']
    base = _item_path(file_id, drive_id)
    data = _graph(f'{base}/workbook/createSession', method='POST',
                  body={'persistChanges': persist})
    sid = data['id']
    _excel_sessions[key] = {'id': sid, 'ts': time.time()}
    return sid

def _graph_excel(path: str, file_id: str, drive_id: str = None,
                 method='GET', body=None, write: bool = False) -> dict:
    """Graph call with workbook session header. Auto-retries once on stale session.
    write=True uses a persist=True session so changes are committed to the file."""
    for attempt in range(2):
        sid = _get_session(file_id, drive_id, persist=write)
        token = get_access_token()
        url = f'https://graph.microsoft.com/v1.0{path}'
        data = json.dumps(body).encode() if body else None
        headers = {'Authorization': f'Bearer {token}',
                   'workbook-session-id': sid}
        if data:
            headers['Content-Type'] = 'application/json'
        req = urllib.request.Request(url, data=data, headers=headers, method=method)
        try:
            with urllib.request.urlopen(req) as r:
                return json.loads(r.read()) if r.length != 0 else {}
        except urllib.error.HTTPError as e:
            if e.code in (404, 409, 410) and attempt == 0:
                _excel_sessions.pop((file_id, write), None)  # invalidate, retry
                continue
            raise

# ---------------------------------------------------------------------------
# Excel — Graph Workbook API
# ---------------------------------------------------------------------------

def _item_path(file_id: str, drive_id: str = None) -> str:
    """Return the correct Graph API path for a drive item (owned or shared)."""
    fid = urllib.parse.quote(file_id, safe='')
    if drive_id:
        did = urllib.parse.quote(drive_id, safe='')
        return f'/drives/{did}/items/{fid}'
    return f'/me/drive/items/{fid}'


def excel_list_sheets(file_id: str, drive_id: str = None) -> list[str]:
    return [s['name'] for s in excel_list_sheets_with_ids(file_id, drive_id)]

def excel_list_sheets_with_ids(file_id: str, drive_id: str = None) -> list[dict]:
    """Return [{id, name}] for all worksheets — used for rename detection."""
    base = _item_path(file_id, drive_id)
    data = _graph_excel(f'{base}/workbook/worksheets?$select=id,name', file_id, drive_id)
    return [{'id': s['id'], 'name': s['name']} for s in data.get('value', [])]

def rebuild_excel_cache(file_id: str, filename: str, drive_id: str = None,
                         file_modified: str = '') -> dict:
    """Rebuild tab descriptions, carrying forward existing descriptions for renamed sheets.
    Returns the new tabs dict {sheet_name: str | list}.
    Tab values may be strings (simple tabs) or lists of section dicts (multi-section tabs).
    Tabs covered by tab_groups need no individual entry — leave them absent from tabs."""
    sheets = excel_list_sheets_with_ids(file_id, drive_id)
    cache = load_excel_cache(file_id)
    old_tabs      = cache.get('tabs', {})       if cache else {}
    old_ids       = cache.get('sheet_ids', {})  if cache else {}
    tab_groups    = cache.get('tab_groups', {}) if cache else {}

    # Build reverse map: graph_id → old_name for rename detection
    id_to_old_name = {v: k for k, v in old_ids.items()}

    # Build set of sheet names covered by tab_groups (no individual description needed)
    grouped = set()
    for g in tab_groups.values():
        pat = g.get('pattern', '')
        if pat:
            for s in sheets:
                if re.match(pat, s['name']):
                    grouped.add(s['name'])

    tabs      = {}
    sheet_ids = {}
    for s in sheets:
        sheet_ids[s['name']] = s['id']
        if s['name'] in grouped:
            continue  # covered by tab_groups, no individual entry needed
        old_name = id_to_old_name.get(s['id'])
        if old_name and old_name in old_tabs:
            tabs[s['name']] = old_tabs[old_name]   # carry forward (str or list)
        elif s['name'] in old_tabs:
            tabs[s['name']] = old_tabs[s['name']]
        else:
            tabs[s['name']] = ''  # new sheet — caller must supply description

    if not file_modified:
        try:
            base = _item_path(file_id, drive_id)
            meta = _graph(f'{base}?$select=lastModifiedDateTime')
            file_modified = meta.get('lastModifiedDateTime', '')
        except Exception:
            pass

    save_excel_cache(file_id, filename, tabs, file_modified=file_modified,
                     sheet_ids=sheet_ids, tab_groups=tab_groups)
    return tabs


def excel_read_range(file_id: str, sheet: str, cell_range: str = 'A1:Z100', drive_id: str = None) -> list[list]:
    """Read a cell range. Returns list of rows, each row a list of values."""
    base = _item_path(file_id, drive_id)
    encoded_sheet = urllib.parse.quote(sheet)
    encoded_range = urllib.parse.quote(cell_range)
    data = _graph_excel(f'{base}/workbook/worksheets/{encoded_sheet}/range(address=\'{encoded_range}\')?$select=values',
                        file_id, drive_id)
    return data.get('values', [])


def excel_write_range(file_id: str, sheet: str, cell_range: str, values: list[list], drive_id: str = None) -> dict:
    """Write values to a cell range. values is a list of rows."""
    base = _item_path(file_id, drive_id)
    encoded_sheet = urllib.parse.quote(sheet)
    encoded_range = urllib.parse.quote(cell_range)
    return _graph_excel(
        f'{base}/workbook/worksheets/{encoded_sheet}/range(address=\'{encoded_range}\')',
        file_id, drive_id, method='PATCH', body={'values': values}, write=True
    )


def excel_read_table(file_id: str, sheet: str, table_name: str, drive_id: str = None) -> list[dict]:
    """Read a named table as list of dicts (column header → value).
    Fetches rows and column headers in parallel."""
    base = _item_path(file_id, drive_id)
    encoded_sheet = urllib.parse.quote(sheet)
    encoded_table = urllib.parse.quote(table_name)
    rows_url    = f'{base}/workbook/worksheets/{encoded_sheet}/tables/{encoded_table}/rows'
    columns_url = f'{base}/workbook/worksheets/{encoded_sheet}/tables/{encoded_table}/columns'
    with ThreadPoolExecutor(max_workers=2) as ex:
        f_rows = ex.submit(_graph_excel, rows_url,    file_id, drive_id)
        f_cols = ex.submit(_graph_excel, columns_url, file_id, drive_id)
        data        = f_rows.result()
        header_data = f_cols.result()
    headers = [c['name'] for c in header_data.get('value', [])]
    return [dict(zip(headers, row.get('values', [[]])[0]))
            for row in data.get('value', [])]


def excel_used_range(file_id: str, sheet: str, drive_id: str = None, max_rows: int = None) -> list[list]:
    """Read the used range of a sheet. Pass max_rows to limit rows returned (faster for large sheets)."""
    base = _item_path(file_id, drive_id)
    encoded_sheet = urllib.parse.quote(sheet)
    data = _graph_excel(f'{base}/workbook/worksheets/{encoded_sheet}/usedRange?$select=values',
                        file_id, drive_id)
    rows = data.get('values', [])
    return rows[:max_rows] if max_rows else rows


def excel_used_range_batch(file_id: str, sheets: list[str], drive_id: str = None,
                            max_rows: int = None, max_workers: int = 8) -> dict[str, list[list]]:
    """Read multiple sheets in parallel. Returns {sheet_name: rows}.
    Failed sheets map to an empty list."""
    def _fetch(sheet):
        try:
            return sheet, excel_used_range(file_id, sheet, drive_id=drive_id, max_rows=max_rows)
        except Exception:
            return sheet, []
    with ThreadPoolExecutor(max_workers=min(max_workers, len(sheets))) as ex:
        return dict(ex.map(_fetch, sheets))


# ---------------------------------------------------------------------------
# Word — download + python-docx
# ---------------------------------------------------------------------------

def word_read(file_id: str, drive_id: str = None) -> str:
    """Download a .docx and return full plain text."""
    from docx import Document
    raw = _graph_bytes(f'{_item_path(file_id, drive_id)}/content')
    doc = Document(io.BytesIO(raw))
    return '\n'.join(p.text for p in doc.paragraphs)


def word_read_structured(file_id: str, drive_id: str = None) -> list[dict]:
    """Return paragraphs with style info: [{'style': ..., 'text': ...}]"""
    from docx import Document
    raw = _graph_bytes(f'{_item_path(file_id, drive_id)}/content')
    doc = Document(io.BytesIO(raw))
    return [{'style': p.style.name, 'text': p.text} for p in doc.paragraphs]


def word_append(file_id: str, paragraphs: list[str], heading: str = None, drive_id: str = None) -> dict:
    """Download, append paragraphs (optionally under a heading), upload back."""
    from docx import Document
    base = _item_path(file_id, drive_id)
    raw = _graph_bytes(f'{base}/content')
    doc = Document(io.BytesIO(raw))
    if heading:
        doc.add_heading(heading, level=2)
    for para in paragraphs:
        doc.add_paragraph(para)
    buf = io.BytesIO()
    doc.save(buf)
    return _graph_upload(
        f'{base}/content',
        buf.getvalue(),
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    )


# ---------------------------------------------------------------------------
# PowerPoint — download + python-pptx
# ---------------------------------------------------------------------------

def pptx_read(file_id: str, drive_id: str = None) -> list[dict]:
    """Download a .pptx and return slide text: [{'slide': N, 'title': ..., 'text': ...}]"""
    from pptx import Presentation
    raw = _graph_bytes(f'{_item_path(file_id, drive_id)}/content')
    prs = Presentation(io.BytesIO(raw))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title = ''
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.shape_type == 13:  # picture
                continue
            text = '\n'.join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
            if shape.name.lower().startswith('title') or (hasattr(slide, 'shapes') and shape == slide.shapes.title):
                title = text
            else:
                if text:
                    texts.append(text)
        slides.append({'slide': i, 'title': title, 'text': '\n'.join(texts)})
    return slides


def pptx_slide_count(file_id: str, drive_id: str = None) -> int:
    """Count slides by reading the ZIP directory — no pptx parsing needed."""
    raw = _graph_bytes(f'{_item_path(file_id, drive_id)}/content')
    with zipfile.ZipFile(io.BytesIO(raw)) as z:
        return sum(1 for n in z.namelist()
                   if n.startswith('ppt/slides/slide') and n.endswith('.xml'))


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main(args):
    if args.cmd == 'search':
        results = search_files(args.query, args.type)
        for f in results:
            print(f"{f['name']} | {f['id']} | {f.get('webUrl', '')}")

    elif args.cmd == 'list-folder':
        items = list_folder(args.path)
        for f in items:
            kind = 'DIR' if 'folder' in f else 'FILE'
            print(f"[{kind}] {f['name']}")

    elif args.cmd == 'excel-sheets':
        fid = args.file_id or get_file_id(args.file_path)
        print('\n'.join(excel_list_sheets(fid)))

    elif args.cmd == 'excel-read':
        fid = args.file_id or get_file_id(args.file_path)
        rows = excel_read_range(fid, args.sheet, args.range)
        for row in rows:
            print('\t'.join(str(c) for c in row))

    elif args.cmd == 'excel-used':
        fid = args.file_id or get_file_id(args.file_path)
        rows = excel_used_range(fid, args.sheet, max_rows=args.max_rows or None)
        for row in rows:
            print('\t'.join(str(c) for c in row))

    elif args.cmd == 'word-read':
        fid = args.file_id or get_file_id(args.file_path)
        text = word_read(fid)
        max_chars = 0 if args.full else args.max_chars
        if max_chars and len(text) > max_chars:
            text = text[:max_chars] + f'\n... [truncated — {len(text)} chars total, use --full for complete text]'
        print(text)

    elif args.cmd == 'pptx-read':
        fid = args.file_id or get_file_id(args.file_path)
        slides = pptx_read(fid)
        max_slides = 0 if args.full else args.max_slides
        if max_slides and len(slides) > max_slides:
            print(f'[{len(slides)} slides total — showing first {max_slides}, use --full for all]')
            slides = slides[:max_slides]
        for slide in slides:
            print(f"--- Slide {slide['slide']}: {slide['title']} ---")
            if slide['text']:
                print(slide['text'])


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Office files CLI via OneDrive')
    sub = parser.add_subparsers(dest='cmd')

    p = sub.add_parser('search', help='Search OneDrive for files')
    p.add_argument('query')
    p.add_argument('--type', help='Filter by extension: xlsx, docx, pptx')

    p = sub.add_parser('list-folder', help='List files in a folder')
    p.add_argument('path', nargs='?', default='/')

    for cmd in ('excel-sheets', 'excel-read', 'excel-used', 'word-read', 'pptx-read'):
        p = sub.add_parser(cmd)
        g = p.add_mutually_exclusive_group(required=True)
        g.add_argument('--file-id', dest='file_id')
        g.add_argument('--file-path', dest='file_path')
        if 'excel-read' in cmd:
            p.add_argument('sheet')
            p.add_argument('range', nargs='?', default='A1:Z200')
        elif 'excel-used' in cmd:
            p.add_argument('sheet', nargs='?')
            p.add_argument('--max-rows', type=int, default=0, dest='max_rows',
                           help='Limit rows returned (default: all). 0 = no limit.')
        elif 'excel' in cmd:
            p.add_argument('sheet', nargs='?')
        elif cmd == 'word-read':
            p.add_argument('--max-chars', type=int, default=8000, dest='max_chars',
                           help='Truncate output at N chars (default 8000). Use --full to disable.')
            p.add_argument('--full', action='store_true', help='Return full content')
        elif cmd == 'pptx-read':
            p.add_argument('--max-slides', type=int, default=20, dest='max_slides',
                           help='Max slides to show (default 20). Use --full to disable.')
            p.add_argument('--full', action='store_true', help='Return all slides')

    args = parser.parse_args()
    if not args.cmd:
        parser.print_help()
    else:
        main(args)
